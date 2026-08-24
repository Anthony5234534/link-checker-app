import pandas as pd
from pptx import Presentation
# 假設你的 PPT 解析函數寫在一個叫 ppt_extractor.py 的檔案中，或者你可以直接放同一個專案
# 如果你是直接寫在主程式，可以把 parse_ppt_to_excel 導入進來
from apify_scraper import run_apify_scraper
import os

def process_ppt_and_scrape(ppt_source) -> pd.DataFrame:
    """
    【內部輔助工具】
    1. 解析 PPT 提取所有超連結與上下文。
    2. 自動呼叫 Apify 爬蟲背景抓取內容與狀態。
    3. 直接整合成最終 DataFrame，供 Streamlit 或後續 AI 比對使用。
    """
    print("step 1: 正在解析 PPT 檔案...")
    # 這裡直接引用你提供的 ppt_parser 邏輯 (假設函數名稱為 parse_ppt_to_excel)
    # 為了方便獨立運作，我們在下方把 parse_ppt_to_excel 的程式碼融合進來或直接調用
    df_ppt = parse_ppt_to_excel_internal(ppt_source)
    
    if df_ppt.empty:
        print("⚠️ 警告：在 PPT 中找不到任何超連結！")
        return df_ppt

    print("step 2: 正在背景執行 Apify 網頁爬蟲...")
    df_final = run_apify_scraper(df_ppt)
    
    return df_final


# 內部引用的 PPT 解析函數 (對應你提供的程式碼)
from urllib.parse import urlparse
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_platform(url):
    try:
        if not url.startswith('http'):
            url = 'https://' + url.replace('https//', 'https://')
        parsed_url = urlparse(url)
        domain = parsed_url.netloc.lower()
        if 'instagram.com' in domain: return 'Instagram'
        elif 'facebook.com' in domain or 'fb.com' in domain: return 'Facebook'
        elif 'threads.net' in domain or 'threads.com' in domain: return 'Threads'
        elif 'xiaohongshu.com' in domain or 'xhslink.com' in domain: return 'Xiaohongshu'
        elif 'douyin.com' in domain: return 'Douyin'
        elif 'weibo.com' in domain or 'weibo.cn' in domain: return 'Weibo'
        elif 'mp.weixin.qq.com' in domain: return 'WeChat Official Account'
        elif 'youtube.com' in domain or 'youtu.be' in domain: return 'YouTube'
        else: return 'Website'
    except Exception:
        return 'Unknown'

def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape

def parse_ppt_to_excel_internal(ppt_source):
    prs = Presentation(ppt_source)
    data = []

    for slide_index, slide in enumerate(prs.slides):
        slide_num = slide_index + 1
        slide_title = "No Title"
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.strip().replace('\n', ' ')

        for shape in iter_shapes(slide.shapes): 
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                runs_info = []
                for run in paragraph.runs:
                    runs_info.append({
                        'text': run.text,
                        'url': run.hyperlink.address if run.hyperlink else None
                    })

                current_preceding = ""
                link_group_context = ""

                for r in runs_info:
                    if r['url']:
                        link_url = r['url']
                        platform = extract_platform(link_url)

                        if current_preceding.strip():
                            cleaned_text = current_preceding.replace('[', '').replace(']', '')
                            link_group_context = cleaned_text.strip().replace('\n', ' ')
                            current_preceding = "" 

                        data.append({
                            'Slide number': slide_num,
                            'Slide_Title': slide_title,
                            'Link_URL': link_url,
                            'Platform': platform,
                            'Preceding_Context': link_group_context
                        })
                    else:
                        current_preceding += r['text']

    df = pd.DataFrame(data)

    if not df.empty:
        df['Preceding_Context'] = df['Preceding_Context'].replace('', pd.NA)
        df['Preceding_Context'] = df.groupby('Slide number')['Preceding_Context'].ffill()
        df['Preceding_Context'] = df['Preceding_Context'].fillna('')

    print(f"PPT 解析成功！共找到 {len(df)} 個連結。")
    return df



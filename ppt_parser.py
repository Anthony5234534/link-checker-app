import re
from pptx.enum.shapes import MSO_SHAPE_TYPE

# For each slide shapes, iterate each shape separately (regardless of whether they are grouped)
# Input the whole slide shapes, output the shape
def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


# Input the ppt file
# Output the yields dict: A dictionary containing details such as slide_num, slide_title, type, url, preceding_context, and associated run references (a unique key for further highlight).
def iter_ppt_links(prs):

    # Iterate each slide in the whole ppt slides
    for slide_index, slide in enumerate(prs.slides):
        slide_num = slide_index + 1
        slide_title = "No Title"

        # Fill in the title if hv
        if slide.shapes.title and slide.shapes.title.text:
            slide_title = slide.shapes.title.text.strip().replace('\n', ' ')

        # Iterate each shape in shapes
        for shape in iter_shapes(slide.shapes): 
            # If table, then should select the whole row or column text    
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)

                # Store corrdinate of all links given a table
                table_links = []
                for r in range(rows):
                    for c in range(cols):
                        cell = table.cell(r, c)
                        for p in cell.text_frame.paragraphs:
                            for run in p.runs:
                                if run.hyperlink and run.hyperlink.address:
                                    table_links.append({
                                        'row': r,
                                        'col': c,
                                        'url': run.hyperlink.address,
                                        'run': run
                                    })
                
                if not table_links:
                    continue 

                # Determine if Horizontal or Vertical
                # if link is mainly in one of the column, then horzontal 
                # if link is mainly in one of the row, then vertical
                
                # Determine if horzontal 
                col_link_counts = {}
                for link in table_links:
                    c = link['col']
                    col_link_counts[c] = col_link_counts.get(c, 0) + 1

                is_horizontal_table = False
                horizontal_col = -1
                expected_links = rows - 1
                tolerance = 2

                # The link should all in one partuiclar column, and the number of link that column should largely mathch to the row number
                for c, count in col_link_counts.items():
                    if count == len(table_links) and count >= (expected_links - tolerance):
                        is_horizontal_table = True
                        horizontal_col = c
                        break

                row_link_counts = {}
                for link in table_links:
                    r = link['row']
                    row_link_counts[r] = row_link_counts.get(r, 0) + 1

                is_vertical_table = False
                vertical_row = -1
                expected_links_v = cols - 1 
                tolerance_v = 2 

                # The link should all in one partuiclar row, and the number of link that row should largely mathch to the column number
                for r, count in row_link_counts.items():
                    if count == len(table_links) and count >= (expected_links_v - tolerance_v):
                        is_vertical_table = True
                        vertical_row = r
                        break

                # if horizontal table, then for that link, extract all row
                if is_horizontal_table:
                    for link_info in table_links:
                        r = link_info['row']
                        link_url = link_info['url']
                        run_obj = link_info['run']

                        row_texts = []
                        for c in range(cols):
                            if c == horizontal_col: 
                                continue 
                            cell_text = table.cell(r, c).text.strip().replace('\n', ' ')
                            if cell_text:
                                row_texts.append(cell_text)
                        
                        preceding_context = " | ".join(row_texts)

                        yield {
                            'slide_num': slide_num,
                            'slide_title': slide_title,
                            'type': 'table_horizontal',
                            'url': link_url,
                            'preceding_context': preceding_context,
                            'run': run_obj
                        }

                elif is_vertical_table:
                    # if vertical table, then for that link, extract all column
                    target_links = [l for l in table_links if l['row'] == vertical_row]
                    
                    for link_info in target_links:
                        c = link_info['col']
                        link_url = link_info['url']
                        run_obj = link_info['run']

                        col_texts = []
                        top_header = table.cell(0, c).text.strip().replace('\n', ' ')
                        if top_header:
                            col_texts.append(top_header)

                        for other_r in range(rows):
                            if other_r == vertical_row: 
                                continue 
                            cell_text = table.cell(other_r, c).text.strip().replace('\n', ' ')
                            row_header = table.cell(other_r, 0).text.strip().replace('\n', ' ')
                            if cell_text and row_header:
                                col_texts.append(f"{row_header}: {cell_text}")
                            elif cell_text:
                                col_texts.append(cell_text)

                        preceding_context = " | ".join(col_texts)

                        yield {
                            'slide_num': slide_num,
                            'slide_title': slide_title,
                            'type': 'table_vertical',
                            'url': link_url,
                            'preceding_context': preceding_context,
                            'run': run_obj
                        }
                # if not vertical or horizontal table, then pass 
                else:
                    continue
                continue 

            # if shape no text, then pass
            if not shape.has_text_frame:
                continue

            # Count the link in each shape    
            actual_links = []
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.hyperlink and run.hyperlink.address:
                        actual_links.append(run)
            
            link_count = len(actual_links)

            # if no. of link = 0, then pass
            if link_count == 0:
                continue

            # if no. of link = 1, then whole text in shape should be extracted.
            if link_count == 1:
                run_obj = actual_links[0]
                link_url = run_obj.hyperlink.address
                
                full_box_text = shape.text_frame.text.strip().replace('\n', ' | ')
                full_box_text = re.sub(r'\[\s*link\s*\]', '', full_box_text, flags=re.IGNORECASE).strip()
                
                yield {
                    'slide_num': slide_num,
                    'slide_title': slide_title,
                    'type': 'single_link',
                    'url': link_url,
                    'preceding_context': full_box_text,
                    'run': run_obj
                }

            # if no. of link > 1, then extract the text before the each link
            else: 
                for paragraph in shape.text_frame.paragraphs:
                    current_preceding = ""
                    current_preceding_runs = []
                    offset = 0
                    
                    for run in paragraph.runs:
                        text_content = run.text
                        link_url = run.hyperlink.address if run.hyperlink else None
                        
                        if link_url:
                            matches = list(re.finditer(r'\[\s*link\s*\]', current_preceding, flags=re.IGNORECASE))
                            last_end = 0
                            
                            for match in matches:
                                context = current_preceding[last_end:match.start()]
                                cleaned_context = re.sub(r'[\[\]]', '', context).strip().replace('\n', ' ')
                                
                                yield {
                                    'slide_num': slide_num,
                                    'slide_title': slide_title,
                                    'type': 'missing_link',
                                    'url': "MISSING_URL_MANUAL_REQUIRED",
                                    'preceding_context': cleaned_context,
                                    'match_start': match.start(),
                                    'match_end': match.end(),
                                    'run_offsets': current_preceding_runs
                                }
                                last_end = match.end()
                            
                            context = current_preceding[last_end:]
                            cleaned_context = re.sub(r'[\[\]]', '', context).strip().replace('\n', ' ')
                            
                            yield {
                                'slide_num': slide_num,
                                'slide_title': slide_title,
                                'type': 'hyperlink',
                                'url': link_url,
                                'preceding_context': cleaned_context,
                                'run': run
                            }
                            
                            current_preceding = ""
                            current_preceding_runs = []
                            offset = 0
                            
                        else:
                            start = offset
                            end = offset + len(text_content)
                            current_preceding_runs.append((run, start, end))
                            current_preceding += text_content
                            offset = end
                
                matches = list(re.finditer(r'\[\s*link\s*\]', current_preceding, flags=re.IGNORECASE))
                last_end = 0
                for match in matches:
                    context = current_preceding[last_end:match.start()]
                    cleaned_context = re.sub(r'[\[\]]', '', context).strip().replace('\n', ' ')
                    
                    yield {
                        'slide_num': slide_num,
                        'slide_title': slide_title,
                        'type': 'missing_link',
                        'url': "MISSING_URL_MANUAL_REQUIRED",
                        'preceding_context': cleaned_context,
                        'match_start': match.start(),
                        'match_end': match.end(),
                        'run_offsets': current_preceding_runs
                    }
                    last_end = match.end()